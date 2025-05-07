from typing import List
from langchain_core.documents import Document
from loguru import logger
from pptx import Presentation
from tqdm import tqdm
from langchain_core.document_loaders import BaseLoader


class PPTLoader(BaseLoader):
    def __init__(self, file_path, load_mode):
        """
        初始化PPT加载器

        Args:
            file_path: PPT文件路径
            load_mode: 加载模式，"full"表示将所有幻灯片合并为一个文档，"page"表示每个幻灯片作为单独文档
        """
        self.file_path = file_path
        self.load_mode = load_mode

    def load(self) -> List[Document]:
        """加载PPT文件并返回文档列表"""
        logger.info(f"解析PPT文件[{self.file_path}]")
        docs = []
        prs = Presentation(self.file_path)

        # 用于全文模式的文本累积器
        full_text = "" if self.load_mode == "full" else None

        for slide_number, slide in tqdm(enumerate(prs.slides, start=1), desc=f"解析[{self.file_path}]的幻灯片"):
            # 用于单页模式的文本累积器
            page_text = "" if self.load_mode == "page" else None

            for shape in slide.shapes:
                # 处理文本框
                if shape.has_text_frame:
                    text = self._extract_text_from_frame(shape.text_frame)
                    if self.load_mode == "full":
                        full_text += text
                    elif self.load_mode == "page":
                        page_text += text

                # 处理表格
                if shape.has_table:
                    table_content = self._extract_text_from_table(shape.table)
                    # 表格总是作为单独文档
                    docs.append(Document(table_content, metadata={"format": "table"}))

            # 对于页面模式，每页内容作为单独文档
            if self.load_mode == "page" and page_text.strip():
                docs.append(Document(page_text.strip(), metadata={"slide_number": slide_number}))

        # 对于全文模式，将所有文本作为一个文档
        if self.load_mode == "full" and full_text.strip():
            docs.append(Document(full_text.strip()))

        return docs

    def _extract_text_from_frame(self, text_frame):
        """从文本框中提取文本"""
        text = ""
        for paragraph in text_frame.paragraphs:
            text += paragraph.text.strip() + "\n"
        return text

    def _extract_text_from_table(self, table):
        """从表格中提取文本"""
        table_content = ""
        for row in table.rows:
            for cell in row.cells:
                table_content += self._extract_text_from_frame(cell.text_frame)
        return table_content.strip()